"""File parsers → list of (page_number, text) tuples.

Supported: pdf, docx, pptx, xlsx, csv, txt, md, html, json, log.
Each parser is deliberately dependency-light so the image stays small.
"""
from __future__ import annotations

import io
import json
from typing import List, Tuple

Pages = List[Tuple[int, str]]

SUPPORTED_EXTENSIONS = {
    "pdf", "docx", "pptx", "xlsx", "csv", "txt", "md", "markdown",
    "html", "htm", "json", "log",
}


def parse_file(filename: str, data: bytes) -> Pages:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "txt"
    if ext == "pdf":
        return _parse_pdf(data)
    if ext == "docx":
        return _parse_docx(data)
    if ext == "pptx":
        return _parse_pptx(data)
    if ext == "xlsx":
        return _parse_xlsx(data)
    if ext in ("html", "htm"):
        return _parse_html(data)
    if ext == "json":
        return _parse_json(data)
    # txt / md / csv / log / anything text-like
    return [(1, data.decode("utf-8", errors="replace"))]


def _parse_pdf(data: bytes) -> Pages:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    pages: Pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append((i, text))
    return pages


def _parse_docx(data: bytes) -> Pages:
    import docx
    doc = docx.Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return [(1, "\n".join(parts))]


def _parse_pptx(data: bytes) -> Pages:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    pages: Pages = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        if texts:
            pages.append((i, "\n".join(texts)))
    return pages


def _parse_xlsx(data: bytes) -> Pages:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    pages: Pages = []
    for i, ws in enumerate(wb.worksheets, start=1):
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            pages.append((i, f"Sheet: {ws.title}\n" + "\n".join(rows)))
    return pages


def _parse_html(data: bytes) -> Pages:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(data, "html.parser")
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    return [(1, soup.get_text(separator="\n", strip=True))]


def _parse_json(data: bytes) -> Pages:
    try:
        obj = json.loads(data)
        return [(1, json.dumps(obj, indent=2, ensure_ascii=False))]
    except Exception:
        return [(1, data.decode("utf-8", errors="replace"))]
